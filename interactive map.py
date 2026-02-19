import ssl
import os
import streamlit as st
import folium
from streamlit_folium import st_folium
from folium.plugins import MeasureControl
from PIL import Image, ImageOps
from PIL.ExifTags import TAGS, GPSTAGS
import base64
from io import BytesIO
import easyocr
import numpy as np
import re
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from google import genai
from google.genai import types
import zipfile
from lxml import etree
import math
from datetime import datetime # เพิ่มโมดูลสำหรับจัดการวันที่

# แก้ไขปัญหา SSL
ssl._create_default_https_context = ssl._create_unverified_context

# --- 1. ตั้งค่า Google Gemini API ---
client = genai.Client(api_key="AIzaSyBHAKfkjkb2wdzAZQZ74dFRD4Ib5Dj6cHY")

@st.cache_resource
def load_ocr():
    model_path = os.path.join(os.getcwd(), "easyocr_models")
    if not os.path.exists(model_path):
        os.makedirs(model_path)
    return easyocr.Reader(['en'], gpu=False, model_storage_directory=model_path)

# โหลด Template พื้นหลัง PowerPoint เก็บไว้ใน Cache
@st.cache_data
def load_template_bytes(file_id):
    try:
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            return response.content
    except: pass
    return None

# --- 2. ฟังก์ชันช่วยดึงรูปภาพ ---
def get_image_base64_from_drive(file_id):
    try:
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            return base64.b64encode(response.content).decode()
    except Exception: return None
    return None

# --- 3. ฟังก์ชันวิเคราะห์สาเหตุด้วย AI ---
def analyze_cable_issue(image_bytes):
    try:
        response = client.models.generate_content(
            model="gemini-1.5-flash",
            contents=[
                """วิเคราะห์รูปภาพสายเคเบิลนี้และเลือกตอบเพียง "หนึ่งเดียว" จาก 4 สาเหตุ:
                1. cable ตกพื้น | 2. หัวต่ออยู่กลาง span เสาไฟฟ้า | 3. ไฟไหม้ cable | 4. หัวต่อขวดน้ำ
                ตอบเฉพาะชื่อสาเหตุภาษาไทยเท่านั้น หากวิเคราะห์ไม่ได้ให้ตอบว่า cable ตกพื้น""",
                types.Part.from_bytes(data=image_bytes, mime_type="image/jpeg")
            ]
        )
        result = response.text.strip()
        if not result or "วิเคราะห์ไม่ได้" in result:
            return "cable ตกพื้น"
        return result
    except Exception:
        return "cable ตกพื้น"

# --- 4. ฟังก์ชันจัดการพิกัด ---
def get_lat_lon_exif(image):
    try:
        exif = image._getexif()
        if not exif: return None, None
        gps_info = {}
        for (idx, tag) in TAGS.items():
            if tag == 'GPSInfo':
                for (t, value) in GPSTAGS.items():
                    if t in exif[idx]: gps_info[value] = exif[idx][t]
        def dms_to_decimal(dms, ref):
            d, m, s = [float(x) for x in dms]
            res = d + (m / 60.0) + (s / 3600.0)
            return -res if ref in ['S', 'W'] else res
        return dms_to_decimal(gps_info['GPSLatitude'], gps_info['GPSLatitudeRef']), \
               dms_to_decimal(gps_info['GPSLongitude'], gps_info['GPSLongitudeRef'])
    except: return None, None

def get_lat_lon_ocr(image):
    try:
        reader = load_ocr() 
        img_for_ocr = image.copy()
        img_for_ocr.thumbnail((1000, 1000)) 
        img_np = np.array(img_for_ocr.convert('RGB'))
        results = reader.readtext(img_np, paragraph=True, allowlist='0123456789.NE ne')
        full_text = " ".join([res[1] for res in results])
        match = re.search(r'(\d+\.\d+)\s*[nN].*?(\d+\.\d+)\s*[eE]', full_text)
        if match: return float(match.group(1)), float(match.group(2))
    except: pass
    return None, None

# --- 5. ฟังก์ชันอ่านไฟล์ KML/KMZ ---
def parse_kml_data(file):
    elements = []
    points_pool = []
    try:
        if file.name.endswith('.kmz'):
            with zipfile.ZipFile(file) as z:
                kml_filename = [n for n in z.namelist() if n.endswith('.kml')][0]
                content = z.read(kml_filename)
        else:
            content = file.getvalue()
        root = etree.fromstring(content)
        ns = {'kml': 'http://www.opengis.net/kml/2.2', 'mwm': 'https://maps.me', 'earth': 'http://earth.google.com/kml/2.2'}
        placemarks = root.xpath('.//kml:Placemark | .//earth:Placemark', namespaces=ns)
        for pm in placemarks:
            name_node = pm.xpath('kml:name/text() | earth:name/text()', namespaces=ns)
            custom_name = pm.xpath('.//mwm:customName/mwm:lang[@code="default"]/text()', namespaces=ns)
            final_name = custom_name[0].strip() if custom_name else (name_node[0].strip() if name_node else "ไม่ระบุชื่อ")
            coords = pm.xpath('.//kml:coordinates/text() | .//earth:coordinates/text()', namespaces=ns)
            if coords:
                pts = [[float(c.split(',')[1]), float(c.split(',')[0])] for c in coords[0].strip().split()]
                elements.append({'name': final_name, 'points': pts, 'is_point': len(pts) == 1})
                for p in pts: points_pool.append(p)
        return elements, points_pool
    except: return [], []

def get_farthest_points(coordinates):
    if not coordinates or len(coordinates) < 2: return None, None
    try:
        if len(coordinates) > 200:
            pts = np.array(coordinates)
            candidates = [pts[pts[:,0].argmax()], pts[pts[:,0].argmin()], pts[pts[:,1].argmax()], pts[pts[:,1].argmin()]]
            test_points = candidates
        else:
            test_points = coordinates
        max_dist = -1
        p1_best, p2_best = None, None
        for i in range(len(test_points)):
            for j in range(i + 1, len(test_points)):
                dist = (test_points[i][0] - test_points[j][0])**2 + (test_points[i][1] - test_points[j][1])**2
                if dist > max_dist:
                    max_dist = dist
                    p1_best, p2_best = test_points[i], test_points[j]
        return p1_best, p2_best
    except: return None, None

def get_osrm_route_head_tail(start_coord, end_coord):
    if not start_coord or not end_coord: return None, 0
    coords_str = f"{start_coord[1]},{start_coord[0]};{end_coord[1]},{end_coord[0]}"
    url = f"http://router.project-osrm.org/route/v1/walking/{coords_str}?overview=full&geometries=geojson"
    try:
        r = requests.get(url, timeout=5)
        if r.status_code == 200:
            data = r.json()
            if "routes" in data and len(data["routes"]) > 0:
                route = data["routes"][0]
                geometry = route["geometry"]["coordinates"]
                distance = route["distance"]
                folium_coords = [[lat, lon] for lon, lat in geometry]
                return folium_coords, distance
    except: pass
    return None, 0

# --- 6. ฟังก์ชันสร้าง Label ชื่อ ---
def create_div_label(name, color="#D9534F"):
    return f'''<div style="font-size: 11px; font-weight: 800; color: {color}; white-space: nowrap; transform: translate(-50%, -150%); background-color: transparent; text-shadow: 2px 2px 4px white, -2px -2px 4px white, 2px -2px 4px white, -2px 2px 4px white; font-family: 'Inter', sans-serif;">{name}</div>'''

def img_to_custom_icon(img, issue_text):
    img_resized = img.copy()
    img_resized.thumbnail((150, 150)) 
    buf = BytesIO()
    img_resized.save(buf, format="JPEG", quality=70)
    img_str = base64.b64encode(buf.getvalue()).decode()
    return f'''
        <div style="position: relative; width: fit-content; background-color: white; padding: 5px; border-radius: 12px; box-shadow: 0px 8px 24px rgba(0,0,0,0.12); border: 2px solid #FF8C42; transform: translate(-50%, -100%); margin-top: -10px;">
            <div style="font-size: 11px; font-weight: 700; color: #2D5A27; margin-bottom: 4px; text-align: center;">{issue_text}</div>
            <img src="data:image/jpeg;base64,{img_str}" style="max-width: 140px; display: block; border-radius: 4px;">
            <div style="position: absolute; bottom: -10px; left: 50%; transform: translateX(-50%); width: 0; height: 0; border-left: 10px solid transparent; border-right: 10px solid transparent; border-top: 10px solid #FF8C42;"></div>
        </div>
    '''

# --- 7. ฟังก์ชันสร้างรายงาน PowerPoint ---
def create_summary_pptx(map_image_bytes, image_list, cable_type, route_distance, issue_kml_elements, template_bytes=None):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(5.625)
    
    # ฟังก์ชันเสริมสำหรับวางภาพ Background Template AIS
    def apply_background(slide):
        if template_bytes:
            slide.shapes.add_picture(BytesIO(template_bytes), 0, 0, width=prs.slide_width, height=prs.slide_height)

    # ==========================================
    # --- หน้าที่ 1: หน้าปก (Cover Slide) ---
    # ==========================================
    slide_cover = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide_cover)
    
    # กล่องข้อความหลัก (จัดให้อยู่โซนสีขาว กว้างประมาณ 7.5 นิ้ว)
    cover_box = slide_cover.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(7.5), Inches(2))
    tf_cover = cover_box.text_frame
    
    # 1. เอกสารประกอบ Imp_NMA-XX
    p_cover1 = tf_cover.paragraphs[0]
    p_cover1.alignment = PP_ALIGN.CENTER
    run1 = p_cover1.add_run()
    run1.text = "เอกสารประกอบ "
    run1.font.size = Pt(32)
    run1.font.color.rgb = RGBColor(0, 86, 179) # สีน้ำเงิน
    
    run2 = p_cover1.add_run()
    run2.text = "Imp_NMA-XX"
    run2.font.size = Pt(36)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(0, 86, 179)
    
    # 2. ข้อมูลนำเสนอปรับปรุง EN-NMA OSP
    p_cover2 = tf_cover.add_paragraph()
    p_cover2.alignment = PP_ALIGN.CENTER
    run3 = p_cover2.add_run()
    run3.text = "ข้อมูลนำเสนอปรับปรุง EN-NMA OSP\n"
    run3.font.size = Pt(28)
    run3.font.color.rgb = RGBColor(0, 86, 179)
    
    # 3. Improve Site XXXX
    p_cover3 = tf_cover.add_paragraph()
    p_cover3.alignment = PP_ALIGN.CENTER
    run4 = p_cover3.add_run()
    run4.text = "Improve Site XXXX"
    run4.font.size = Pt(36)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0, 86, 179)

    # 4. กล่องข้อความมุมล่างซ้าย: ดึงวันที่ปัจจุบันมาแสดง
    ver_box = slide_cover.shapes.add_textbox(Inches(0.2), Inches(5.1), Inches(4), Inches(0.5))
    p_ver = ver_box.text_frame.paragraphs[0]
    current_date_str = datetime.now().strftime("%d/%m/%Y") # จัดรูปแบบเป็น วัน/เดือน/ปี
    p_ver.text = f"Ver.Update  {current_date_str}"
    p_ver.font.size = Pt(12)
    p_ver.font.color.rgb = RGBColor(0, 0, 0) # สีดำ

    # ==========================================
    # --- หน้าที่ 2: รายละเอียดสรุป ---
    # ==========================================
    slide0 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide0) 
    
    title_box = slide0.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(7.5), Inches(1))
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = f"รายงานสรุปแนวทางแก้ไขปัญหาและเสนอคร่อม Cable ({cable_type} Core)"
    p_title.font.bold = True
    p_title.font.size = Pt(22)
    
    info_box = slide0.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.5), Inches(3.5))
    tf = info_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]; p1.text = f"• Type Cable: {cable_type} Core"; p1.font.size = Pt(16)
    p2 = tf.add_paragraph(); p2.text = f"• ระยะคร่อม Cable รวม: {route_distance:,.0f} เมตร ({route_distance/1000:.3f} กม.)"; p2.font.size = Pt(16)
    p3 = tf.add_paragraph(); p3.text = f"• รายละเอียดจุดปัญหา:"; p3.font.bold = True; p3.font.size = Pt(16)
    
    for el in issue_kml_elements[:10]:
        p_el = tf.add_paragraph()
        p_el.text = f"  - {el['name']} (Lat: {el['points'][0][0]:.5f}, Long: {el['points'][0][1]:.5f})"
        p_el.font.size = Pt(12)

    # ==========================================
    # --- หน้าที่ 3: ภาพแสดงแผนที่ ---
    # ==========================================
    if map_image_bytes:
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide1)
        
        # แสดงรูปภาพแผนที่เต็มจอ
        slide1.shapes.add_picture(BytesIO(map_image_bytes), 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Title: Topology Overall (ขีดเส้นใต้ มุมบนซ้าย)
        title_box1 = slide1.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(5), Inches(0.5))
        p_title1 = title_box1.text_frame.paragraphs[0]
        p_title1.text = "Topology Overall"
        p_title1.font.bold = True
        p_title1.font.size = Pt(24)
        p_title1.font.underline = True
        
    # ==========================================
    # --- หน้าที่ 4: รูปภาพแสดงจุดที่มีปัญหา ---
    # ==========================================
    if image_list:
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide2)
        
        # Title: รูปภาพแสดงจุดที่มีปัญหา (ขีดเส้นใต้ มุมบนซ้าย)
        title_box2 = slide2.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(6), Inches(0.5))
        p_title2 = title_box2.text_frame.paragraphs[0]
        p_title2.text = "รูปภาพแสดงจุดที่มีปัญหา"
        p_title2.font.bold = True
        p_title2.font.size = Pt(22)
        p_title2.font.underline = True

        cols, rows = 4, 2
        img_w, img_h = Inches(1.8), Inches(1.3)
        margin_x = (Inches(7.8) - (img_w * cols)) / (cols + 1)
        margin_y = Inches(0.8) 
        
        for i, item in enumerate(image_list[:8]):
            curr_row, curr_col = i // cols, i % cols
            x, y = margin_x + (curr_col * (img_w + margin_x)), margin_y + (curr_row * (img_h + Inches(0.8))) 
            image = item['img_obj'].copy()
            target_ratio = img_w / img_h
            w_px, h_px = image.size
            if (w_px/h_px) > target_ratio:
                new_w = h_px * target_ratio
                left = (w_px - new_w) / 2
                image = image.crop((left, 0, left + new_w, h_px))
            else:
                new_h = w_px / target_ratio
                top = (h_px - new_h) / 2
                image = image.crop((0, top, w_px, top + new_h))
            buf = BytesIO(); image.save(buf, format="JPEG"); buf.seek(0)
            slide2.shapes.add_picture(buf, x, y, width=img_w, height=img_h)
            txt_box = slide2.shapes.add_textbox(x, y + img_h + Inches(0.02), img_w, Inches(0.6))
            tf_img = txt_box.text_frame
            tf_img.word_wrap = True
            p1_img = tf_img.paragraphs[0]; p1_img.text = f"สาเหตุ: {item['issue']}"; p1_img.font.size = Pt(8); p1_img.font.bold = True
            p2_img = tf_img.add_paragraph(); p2_img.text = f"Lat: {item['lat']:.5f}\nLong: {item['lon']:.5f}"; p2_img.font.size = Pt(7)
            
    output = BytesIO(); prs.save(output)
    return output.getvalue()

# --- 8. UI Layout ---
st.set_page_config(page_title="AI Cable Survey", layout="wide")
st.markdown("""<style>
    .stApp { background: linear-gradient(120deg, #FFF5ED 0%, #F0F9F1 100%); }
    .header-container { display: flex; align-items: center; justify-content: space-between; padding: 25px; background: white; border-radius: 24px; border-bottom: 5px solid #FF8C42; margin-bottom: 30px; }
    .main-title { background: linear-gradient(90deg, #2D5A27 0%, #FF8C42 100%); -webkit-background-clip: text; -webkit-text-fill-color: transparent; font-weight: 800; font-size: 2.6rem; margin: 0; }
    .joker-icon { width: 100px; height: 100px; object-fit: cover; border-radius: 50%; border: 4px solid #FFFFFF; outline: 3px solid #FF8C42; }
    .stButton>button { background: #2D5A27; color: white; border-radius: 14px; padding: 12px 35px; font-weight: 600; width: 100%; }
</style>""", unsafe_allow_html=True)

joker_base64 = get_image_base64_from_drive("1_G_r4yKyBA_vv3Nf8SdFpQ8UKv4bPLBr")
header_html = f'''<div class="header-container"><div><h1 class="main-title">AI Cable Plotter</h1><p style="margin:0; color: #718096; font-weight: 600;">By Joker EN-NMA</p></div>{"<img src='data:image/png;base64,"+joker_base64+"' class='joker-icon'>" if joker_base64 else ""}</div>'''
st.markdown(header_html, unsafe_allow_html=True)

# --- 9. เมนู KML/KMZ ---
st.subheader("🌐 1. ข้อมูลโครงข่าย & จุดติดตั้ง (KML/KMZ)")
kml_file_yellow = st.file_uploader("Import KMZ - Overall (ภาพรวมแผนที่)", type=['kml', 'kmz'])
kml_file = st.file_uploader("Import KMZ - พิกัดที่มีปัญหาและเสนอคร่อม cable", type=['kml', 'kmz'])

zoom_bounds = []
kml_elements, kml_points_pool, yellow_elements = [], [], []

if kml_file_yellow:
    yellow_elements, _ = parse_kml_data(kml_file_yellow)
    for el in yellow_elements: zoom_bounds.extend(el['points'])
if kml_file:
    kml_elements, kml_points_pool = parse_kml_data(kml_file)
    for el in kml_elements: zoom_bounds.extend(el['points'])

st.markdown("<hr>", unsafe_allow_html=True)

# --- 10. ส่วนรูปภาพสำรวจ ---
st.subheader("📁 2. อัปโหลดรูปภาพสำรวจ")
uploaded_files = st.file_uploader("ลากและวางไฟล์ที่นี่", type=['jpg','jpeg','png'], accept_multiple_files=True, key="survey_uploader")
if 'export_data' not in st.session_state: st.session_state.export_data = []

if uploaded_files:
    current_hash = "".join([f.name + str(f.size) for f in uploaded_files])
    if 'last_hash' not in st.session_state or st.session_state.last_hash != current_hash:
        st.session_state.export_data, st.session_state.last_hash = [], current_hash
    for i, f in enumerate(uploaded_files):
        if i >= len(st.session_state.export_data):
            raw_data = f.getvalue()
            raw_img = Image.open(BytesIO(raw_data))
            img_st = ImageOps.exif_transpose(raw_img)
            lat, lon = get_lat_lon_exif(raw_img)
            if lat is None: lat, lon = get_lat_lon_ocr(img_st)
            if lat:
                issue = analyze_cable_issue(raw_data)
                storage_img = img_st.copy()
                storage_img.thumbnail((1200, 1200))
                st.session_state.export_data.append({'img_obj': storage_img, 'issue': issue, 'lat': lat, 'lon': lon})

for data in st.session_state.export_data: zoom_bounds.append([data['lat'], data['lon']])

route_coords, route_distance = None, 0
if kml_points_pool:
    try:
        f_p = get_farthest_points(kml_points_pool)
        if f_p and f_p[0] is not None and f_p[1] is not None:
            route_coords, route_distance = get_osrm_route_head_tail(f_p[0], f_p[1])
    except: pass

if uploaded_files or kml_elements or yellow_elements:
    m = folium.Map(location=[13.75, 100.5], zoom_start=17, tiles=None, control_scale=True)
    folium.TileLayer(tiles="https://mt1.google.com/vt/lyrs=m&x={x}&y={y}&z={z}", attr="Google", name="Google Maps", opacity=0.4, overlay=False).add_to(m)
    if route_coords:
        folium.PolyLine(route_coords, color="#FF0000", weight=5, opacity=0.8, dash_array='10, 10').add_to(m)
        st.info(f"📍 ระยะคร่อม cable: {route_distance/1000:.3f} กม. ({route_distance:,.0f} เมตร)")
    for elem in yellow_elements:
        if elem['is_point']:
            folium.Marker(elem['points'][0], icon=folium.Icon(color='orange')).add_to(m)
            folium.Marker(elem['points'][0], icon=folium.DivIcon(html=create_div_label(elem['name'], "#CC9900"))).add_to(m)
        else: folium.PolyLine(elem['points'], color="#FFD700", weight=4, opacity=0.8).add_to(m)
    for elem in kml_elements:
        if elem['is_point']:
            folium.Marker(elem['points'][0], icon=folium.Icon(color='red')).add_to(m)
            folium.Marker(elem['points'][0], icon=folium.DivIcon(html=create_div_label(elem['name'], "#D9534F"))).add_to(m)
        else: folium.PolyLine(elem['points'], color="gray", weight=2, opacity=0.4, dash_array='5').add_to(m)
    for d in st.session_state.export_data: folium.Marker([d['lat'], d['lon']], icon=folium.DivIcon(html=img_to_custom_icon(d['img_obj'], d['issue']))).add_to(m)
    m.add_child(MeasureControl(position='topright', primary_length_unit='meters'))
    if zoom_bounds: m.fit_bounds(zoom_bounds, padding=[50, 50])
    st_folium(m, height=1200, use_container_width=True, key="survey_map")

st.markdown("<hr>", unsafe_allow_html=True)
st.subheader("📄 3. สร้างรายงาน PowerPoint")
col_c1, col_c2 = st.columns(2)
with col_c1:
    cable_type = st.selectbox("เลือก Type Cable", ["4", "6", "12", "24", "48", "96"])
    map_cap = st.file_uploader("อัปโหลดรูป Capture แผนที่", type=['jpg','png'])
if map_cap:
    with col_c2:
        if st.button("🚀 ดาวน์โหลดรายงาน PPTX"):
            try:
                # โหลด Template รูปภาพ AIS จากลิงก์ Google Drive
                bg_template_id = "1EqtiR6CVnsbsVIg5Gk5j1v901YXYzjkz"
                template_bytes = load_template_bytes(bg_template_id)
                
                pptx_data = create_summary_pptx(
                    map_cap.getvalue(), 
                    st.session_state.export_data, 
                    cable_type, 
                    route_distance, 
                    kml_elements, 
                    template_bytes # ส่ง Template เข้าไป
                )
                st.download_button(
                    "📥 คลิกเพื่อดาวน์โหลด", 
                    data=pptx_data, 
                    file_name=f"Cable_Survey_{cable_type}C.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"เกิดข้อผิดพลาดในการสร้างรายงาน: {e}")
